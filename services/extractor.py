import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import logging
import easyocr
import numpy as np
from PIL import Image

logger = logging.getLogger(__name__)

# Initialize EasyOCR reader (will download models on first run if not present)
# Using english only for general medical text OCR
ocr_reader = easyocr.Reader(['en'], gpu=False)

def extract_text_from_image(img_bytes: bytes) -> str:
    try:
        # EasyOCR works best with numpy arrays
        image = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        image_np = np.array(image)
        
        # Read text
        results = ocr_reader.readtext(image_np)
        
        # Join extracted text snippets
        text = " ".join([res[1] for res in results])
        return text
    except Exception as e:
        logger.error(f"Error extracting Image OCR: {e}")
        return ""

def extract_text_from_pdf(contents: bytes) -> str:
    try:
        doc = fitz.open(stream=contents, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return ""

def extract_text_from_docx(contents: bytes) -> str:
    try:
        doc = docx.Document(io.BytesIO(contents))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return ""

def extract_text_from_pptx(contents: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(contents))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return ""
